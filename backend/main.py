import sys
import io
import queue
import threading
import asyncio
import json
import re
import os
from fastapi import FastAPI, Request, UploadFile, File, Form, HTTPException
from fastapi.responses import StreamingResponse
from fastapi.middleware.cors import CORSMiddleware
from pydantic import BaseModel
from typing import Optional
from dotenv import load_dotenv

# Document parsers
import pypdf
from docx import Document
from pptx import Presentation

from src.crew import SmartStudyCrew
from src.config.settings import settings

load_dotenv()

# Use Key 2 as primary because Key 1 is known to be exhausted until Feb 13
api_key_primary = os.getenv("GOOGLE_API_KEY_2") or os.getenv("GOOGLE_API_KEY")

# CRITICAL: Disable OpenAI and LiteLLM noise
os.environ["OPENAI_API_KEY"] = "none"
os.environ["LITELLM_LOGGING"] = "False"
os.environ["LITELLM_MODE"] = "native"
os.environ["OTEL_SDK_DISABLED"] = "true"


if api_key_primary:
    os.environ["GEMINI_API_KEY"] = api_key_primary
    os.environ["GOOGLE_API_KEY"] = api_key_primary

# Disable all tracking and telemetry
os.environ["CREWAI_SKIP_TELEMETRY"] = "true"
os.environ["CREWAI_TRACING_ENABLED"] = "false"
os.environ["LANGCHAIN_TRACING_V2"] = "false"


app = FastAPI()


app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class StudyRequest(BaseModel):
    topic: str
    notes: str = ""

# Regex to strip ANSI color codes for clean frontend logs
ANSI_ESCAPE = re.compile(r'\x1B(?:[@-Z\\-_]|\[[0-?]*[ -/]*[@-~])')

class QueueWriter(io.TextIOBase):
    def __init__(self, q):
        self.q = q
    def write(self, data):
        if data:
            clean_data = ANSI_ESCAPE.sub('', data)
            if clean_data.strip():
                self.q.put(clean_data)

@app.post("/upload")
async def upload_file(file: UploadFile = File(...)):
    """Handles document uploads (PDF, DOCS, PPTX, TXT, MD) and extracts text"""
    try:
        content = await file.read()
        filename = file.filename.lower()
        
        # Save file to materials directory
        save_path = settings.output_dir / "materials" / filename
        with open(save_path, "wb") as f:
            f.write(content)
            
        extracted_text = ""

        if filename.endswith('.pdf'):
            try:
                reader = pypdf.PdfReader(io.BytesIO(content))
                for page in reader.pages:
                    extracted_text += page.extract_text() + "\n"
            except Exception as e:
                print(f"PDF extract error: {e}")
        
        elif filename.endswith('.docx') or filename.endswith('.doc'):
            try:
                doc = Document(io.BytesIO(content))
                for para in doc.paragraphs:
                    extracted_text += para.text + "\n"
            except Exception as e:
                print(f"DOCX extract error: {e}")
        
        elif filename.endswith('.pptx'):
            try:
                prs = Presentation(io.BytesIO(content))
                for slide in prs.slides:
                    for shape in slide.shapes:
                        if hasattr(shape, "text"):
                            extracted_text += shape.text + "\n"
            except Exception as e:
                 print(f"PPTX extract error: {e}")
        
        elif filename.endswith('.txt') or filename.endswith('.md'):
            extracted_text = content.decode('utf-8', errors='ignore')
            
        else:
             # Fallback try text decode
            extracted_text = content.decode('utf-8', errors='ignore')

        return {"text": extracted_text.strip(), "filename": file.filename, "saved_path": str(save_path)}

    except Exception as e:
        print(f"Upload error: {e}")
        return {"error": f"Failed to process file: {str(e)}", "text": ""}

@app.post("/generate-plan")
async def generate_plan(request: StudyRequest):
    log_queue = queue.Queue()
    result_queue = queue.Queue()
    
    def run_crew():
        sys.stdout = QueueWriter(log_queue)
        try:
            crew = SmartStudyCrew(request.topic, request.notes)
            result, memory = crew.run()
            
            # Send final report with memory summary
            report_text = str(result)
            memory_summary = memory.get_context_summary()
            
            print(f"\n[FINAL_REPORT]\n{report_text}")
            print(f"\n[MEMORY_SUMMARY]\n{memory_summary}")
            
        except Exception as e:
            print(f"Error during mission: {str(e)}")
        finally:
            sys.stdout = sys.__stdout__
            log_queue.put(None)

    threading.Thread(target=run_crew).start()

    async def stream_generator():
        while True:
            try:
                chunk = await asyncio.to_thread(log_queue.get, timeout=2.0)
                if chunk is None: break
                yield chunk
            except queue.Empty:
                await asyncio.sleep(0.1)
                continue

    return StreamingResponse(stream_generator(), media_type="text/plain")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="127.0.0.1", port=8081)
